from datetime import datetime, timedelta, timezone
from email.mime.image import MIMEImage
import io
from fastapi import BackgroundTasks
import json
from openai import OpenAI
import os
from dotenv import load_dotenv
import re
from sqlite3 import IntegrityError
import string
from typing import Counter, Dict, List, Literal, Optional
import unicodedata
from uuid import uuid4
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile, Depends, status
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
import random
from fastapi.encoders import jsonable_encoder
from fastapi.responses import HTMLResponse, StreamingResponse
import nltk
import spacy
from sqlalchemy import UUID
import urllib
from models import SessionLocal, User, Document, Category
from sqlalchemy.orm import Session, Query
from passlib.context import CryptContext
from jose import JWTError, jwt
from baseclass import BaseClass
from models import VerificationRun, Verification, Proof, EmailProof, VerificationRunDuplicate, Quiz, Question, Answer, QuizResult, ModerationLog
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from urllib.parse import quote

from fastapi.middleware.cors import CORSMiddleware

from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse
import boto3
import textract
from io import BytesIO
from Questgen import main
from deep_translator import GoogleTranslator
import difflib
import nltk
from nltk.corpus import wordnet
from random import shuffle


from pydantic import BaseModel, EmailStr
from mangum import Mangum
import boto3
from io import BytesIO
import uuid


AWS_ACCESS_KEY_ID = os.getenv('AWS_ACCESS_KEY_ID')
AWS_SECRET_ACCESS_KEY = os.getenv('AWS_SECRET_ACCESS_KEY')
AWS_REGION_NAME = os.getenv('AWS_REGION_NAME')
S3_BUCKET_NAME = os.getenv('S3_BUCKET_NAME')

AWS_ACCESS_KEY_ID = os.getenv('AWS_ACCESS_KEY_ID')
AWS_SECRET_ACCESS_KEY = os.getenv('AWS_SECRET_ACCESS_KEY')
AWS_REGION_NAME = os.getenv('AWS_REGION_NAME')
S3_BUCKET_NAME = os.getenv('S3_BUCKET_NAME')

class Book(BaseModel):
    name: str
    genre: Literal["fiction", "non-fiction"]
    price: float
    book_id: Optional[str] = uuid4().hex


BOOKS_FILE = "books.json"
BOOKS = []

if os.path.exists(BOOKS_FILE):
    with open(BOOKS_FILE, "r") as f:
        BOOKS = json.load(f)
        
SECRET_KEY = "YOUR_SECRET_KEY"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 120


app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Ezt módosítsd biztonságosabbra, ha kell
    allow_credentials=True,
    allow_methods=["*"],  # Engedélyez minden metódust
    allow_headers=["*"],  # Engedélyez minden fejlécet
)

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")



app.mount("/static", StaticFiles(directory="static"), name="static")
app.mount("/quiz", StaticFiles(directory="quiz"), name="quiz")
app.mount("/catalog", StaticFiles(directory="catalog"), name="catalog")
app.mount("/main", StaticFiles(directory="main"), name="main")
app.mount("/moderation", StaticFiles(directory="moderation"), name="moderation")
app.mount("/trending", StaticFiles(directory="trending"), name="trending")
app.mount("/quizzes", StaticFiles(directory="quizzes"), name="quizzes")
app.mount("/allquizzes", StaticFiles(directory="allquizzes"), name="allquizzes")

handler = Mangum(app) 
# A token generálása
def create_access_token(data: dict, expires_delta: timedelta = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)):
    to_encode = data.copy()
    expire = datetime.utcnow() + expires_delta
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt
# Dependency
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

# Token validálása
def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Invalid token",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise credentials_exception
        user = get_user_from_db(email, db)
        if user is None:
            raise credentials_exception
        return user
    except JWTError:
        raise credentials_exception
    
def get_user_from_db(email: str, db: Session = SessionLocal()) -> Optional[User]:
    user = db.query(User).filter(User.email == email).first()
    db.close()
    return user

def verify_get_user_from_db(email: str, db: Session = SessionLocal()) -> Optional[User]:
    user = db.query(User).filter(User.email == email).first()
    user.verified = True

    if email.endswith("@inf.elte.hu"):
        user.tokens = 6  
    else:
        user.tokens = 3  

    db.commit()
    db.close()

    return True


def verify_password(plain_password, password_hash):
    return pwd_context.verify(plain_password, password_hash)

@app.post("/token")
async def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends()):
    user = get_user_from_db(form_data.username)
    print("debug")
    print(user)
    if user is None or not verify_password(form_data.password, user.password_hash):
        print("DDDDDDDDDDDDDDDEBUG")
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid credentials"+form_data.username + form_data.password + user.password_hash,
        )
    print("token generating")
    if not user.verified:
        return {"id": user.id, "status": "not_verified", "message": "Email not verified. Please enter the verification code sent to your email."}
    
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={"sub": form_data.username}, expires_delta=access_token_expires
    )
    print("ACCESS TOKEN" + access_token)

    return {"access_token": access_token, "token_type": "bearer"}




# Pydantic models
class UserCreate(BaseModel):
    name: str
    email: EmailStr
    password: str

class UserResponse(BaseModel):
    id: str
    name: str
    email: EmailStr
    role: str

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

def hash_password(password: str):
    return pwd_context.hash(password)

@app.post("/users/", response_model=UserResponse)
def create_user(user: UserCreate, db: Session = Depends(get_db)):
    try:
        if db.query(User).filter(User.email == user.email).first():
            raise HTTPException(status_code=400, detail="Ez az email cím már regisztrálva van!")
        
        if db.query(User).filter(User.name == user.name).first():
            raise HTTPException(status_code=400, detail="Ez a név már foglalt!")

        db_user = User(
            id=str(uuid.uuid4()), 
            name=user.name, 
            email=user.email, 
            password_hash=hash_password(user.password), 
            role="user"
        )


        db.add(db_user)
        db.commit()
        db.refresh(db_user)

            # Ha email már létezik, akkor dobj hibát


        base = BaseClass()

        run_duplicate=base.is_run_duplicate(entity_id=db_user.id, verification_process="EMAIL", session=db)

        if run_duplicate!="":
            new_duplicate_run = VerificationRunDuplicate(
                id=f"verification_verificationrunduplicate_{uuid.uuid4()}",
                serviceProviderID="VB",
                verificationTypeCode="EMAIL",
                entityType=db_user.role,
                entityID=db_user.id,
                verificationProcessCode="EMAIL",
                originalVerificationRunID=run_duplicate
            )
            base.create_verification_run_duplicate(new_duplicate_run)
            return {"status": "DUPLICATE_RUN_FOUND"}
        
        VERIFICATION_EXPIRE_DAYS=5000
        CODE_LENGTH=6
        TRY_EXPIRE_HOURS=24
        MAX_RRETRY_PROCESS=3
        MAX_RETRY_PROCESS_WAIT_TIME_MINUTES=3
        MAX_RETRY_PROCESS_METHOD="EXPONENTIAL"
        MAX_RETRY=3

        new_run = VerificationRun(
                id=f"verification_verificationrun_{uuid.uuid4()}",
                serviceProviderID="VB",
                verificationProcessCode="EMAIL",
                entityType=db_user.role,
                entityID=db_user.id,
                verificationTypeCode="EMAIL",
                status="ONGOING",
                vendor_status="PENDING",
                fail_reason="",
                try_count=0,
                effective_date=datetime.now(),
                expiration_date=datetime.now() + timedelta(hours=TRY_EXPIRE_HOURS),
                remaining_tries=MAX_RETRY
            )

        created_run = base.create_verification_run(new_run, session=db)

        prefix = ''.join([str(random.randint(0, 9)) for _ in range(3)])

        verification_code = ''.join([str(random.randint(0, 9)) for _ in range(CODE_LENGTH)])

        new_proof = EmailProof(
            id=f"verification_proof_{uuid.uuid4()}",
            verificationRunID=created_run.id,
            main_param=db_user.email,
            verification_code=verification_code,
            uploadDate=datetime.now(),
            expirationDate=datetime.now() + timedelta(hours=TRY_EXPIRE_HOURS),
            entityType=db_user.role,
            entityID=db_user.id,
            prefix=prefix,
            ip_address="",
            correct_code_submission_time=None,
            status="PENDING"
        )
        created_proof=base.create_proof(new_proof, session=db)

        proof = created_run.proofs[0]

        phoneresult=base.email_duplicate_check(db_user.id, email=db_user.id, session=db)
        if phoneresult=="":
            #base.update_proof_status(proof=proof, new_status="PENDING", main_param=normalized_phone, session=session)
            proof.status="PENDING"
            proof.main_param=db_user.email

        #send email to address
        #send_email(db_user.email, verification_code)
        verification_code = verification_code
        send_email(db_user.email, verification_code, db_user.name)
        "EMAIL SENT"
    except IntegrityError as e:
        raise HTTPException(status_code=400, detail="Ez az email cím már regisztrálvaA van!")
    return db_user

@app.get("/users/{user_id}", response_model=UserResponse)
def get_user(user_id: str, db: Session = Depends(get_db)):
    db_user = db.query(User).filter(User.id == user_id).first()
    if db_user is None:
        raise HTTPException(status_code=404, detail="User not found")
    return db_user


@app.get("/users", response_model=List[UserResponse])
async def get_users(db: Session = Depends(get_db)):
    db_users = db.query(User).all()  # Query to get all users
    return db_users

# Alap oldal
@app.get("/", response_class=HTMLResponse)
async def index():
    return open("static/index.html").read()


def sanitize_filename(filename: str) -> str:
    """
    Megtisztítja a fájlnevet:
    - Eltávolítja az érvénytelen karaktereket
    - Lecseréli az ékezetes karaktereket (pl. ő -> o, é -> e)
    - Megőrzi az eredeti fájlkiterjesztést
    """
    # Fájlkiterjesztés megőrzése
    name, ext = filename.rsplit('.', 1) if '.' in filename else (filename, '')

    # Ékezetes karakterek normalizálása
    name = unicodedata.normalize('NFKD', name).encode('ascii', 'ignore').decode('ascii')

    # Csak engedélyezett karakterek: betűk, számok, kötőjel és aláhúzás
    name = re.sub(r'[^a-zA-Z0-9_\-]', '_', name)

    # Ha nincs kiterjesztés, nem kell pont
    return f"{name}.{ext}" if ext else name

s3 = boto3.client('s3', aws_access_key_id=AWS_ACCESS_KEY_ID, aws_secret_access_key=AWS_SECRET_ACCESS_KEY, region_name=AWS_REGION_NAME)

@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...), 
    title: str = Form(""), 
    description: str = Form(""), 
    category_id: str = Form(""),
    role: str = Form(""), 
    uploaded_by: str = Form(""), 
    is_edit: bool = Form(False), 
    db: Session = Depends(get_db)
) -> Dict[str, str]:
    """
    Feltölti a fájlt az S3-ba és elmenti a dokumentum adatait az adatbázisba.


    
    """

    MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 MB

    # Fájlméret ellenőrzése
    file.file.seek(0, 2)  # Seek to end to get size
    file_size = file.file.tell()
    file.file.seek(0)  # Seek back to start
    file_size_in_mb = file_size / (1024 * 1024)

    if file_size > MAX_FILE_SIZE:
        return {
            "message": "ERROR",
            "error": str(file_size_in_mb)
        }
    

    print(f"Category ID: {category_id}, Uploaded By: {uploaded_by}")

    category = db.query(Category).filter(Category.id == category_id).first()

    categoryName = category.name
    randomize_it = ''.join(random.choices(string.ascii_lowercase + string.digits, k=6))
    sanitized_category_name = sanitize_filename(categoryName)
    sanitized_filename = sanitize_filename(file.filename)

    filenameNew = f"{randomize_it}_{sanitized_category_name}_{sanitized_filename}"
    # Fájl feltöltése az S3-ba
    s3.upload_fileobj(file.file, S3_BUCKET_NAME, filenameNew)
    file_url = f"https://{S3_BUCKET_NAME}.s3.{AWS_REGION_NAME}.amazonaws.com/{filenameNew}"

    # Új dokumentum mentése az adatbázisba
    if role == 'user':
        new_document = Document(
            title=title,
            description=description,
            file_path=file_url,
            uploaded_by=uploaded_by,
            status="pending",  
            category_id=category_id,
            is_edit=is_edit
        )
        db.add(new_document)
        db.commit()
        db.refresh(new_document)

        return {
            "message": "File is uploaded successfully, and is waiting for approval.",
            "file_url": file_url,
            "document_id": new_document.id,
            "title": new_document.title,
            "description": new_document.description,
            "uploaded_by": new_document.uploaded_by,
            "uploaded_at": new_document.uploaded_at.isoformat(),
        }
    else:
        new_document = Document(
            title=title,
            description=description,
            file_path=file_url,
            uploaded_by=uploaded_by,
            status="approved",  
            category_id=category_id,
        )
        db.add(new_document)
        db.commit()
        db.refresh(new_document)

        return {
            "message": 'File is uploaded successfully.',
            "file_url": file_url,
            "document_id": new_document.id,
            "title": new_document.title,
            "description": new_document.description,
            "uploaded_by": new_document.uploaded_by,
            "uploaded_at": new_document.uploaded_at.isoformat(),
        }
    


@app.delete("/delete/{filename}")
async def delete_file(
    filename: str= "", 
    current_user: User = Depends(get_current_user),  # Aktuális felhasználó lekérése
    db: Session = Depends(get_db)
) -> Dict[str, str]:
    """
    Törli a fájlt az S3-ból, ha a felhasználó admin, moderátor, vagy a fájl feltöltője.
    """
    print("DEBUG: Deleting file...")
    try:
        # Dokumentum adatainak lekérése a filename alapján
        file_path="https://poseidonb.s3.eu-north-1.amazonaws.com/"+filename
        document = db.query(Document).filter(Document.file_path == file_path).first()
        
        if document is None:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not founddd.")
        
        # Jogosultságok ellenőrzése
        if current_user.role not in ["admin", "moderator"] and current_user.id != document.uploaded_by:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN, 
                detail="You do not have permission to delete this file."
            )

        # Fájl törlése az S3-ból
        response = s3.delete_object(Bucket=S3_BUCKET_NAME, Key=filename)
        print(f"DEBUG: S3 delete response: {response}")
        db.delete(document)
        db.commit()
        # Válasz visszaadása
        return {"message": f"File {filename} deleted successfully."}
    
    except Exception as e:
        return {"message": f"Error deleting file: {str(e)}"}

    



@app.get("/download/{filename}")
async def download_file(filename: str):
    s3 = boto3.client('s3', aws_access_key_id=AWS_ACCESS_KEY_ID, aws_secret_access_key=AWS_SECRET_ACCESS_KEY, region_name=AWS_REGION_NAME)

    try:
        # Fájl lekérése az S3-ból
        file_obj = s3.get_object(Bucket=S3_BUCKET_NAME, Key=filename)
        
        # Fájl tartalmának betöltése
        file_content = file_obj['Body'].read()
        
        # Automatikus letöltés (Content-Disposition header)
        return StreamingResponse(BytesIO(file_content), media_type="application/octet-stream", headers={"Content-Disposition": f"attachment; filename={filename}"})
    
    except Exception as e:
        return {"error": str(e)}
    

@app.get("/me", response_model=UserResponse)
def get_current_user_info(user: User = Depends(get_current_user)):
    return user

@app.post("/expire_ongoing_verification_runs")
def expire_ongoing_verification_runs(
    db: Session = Depends(get_db)

):
    base = BaseClass()
    base.expire_ongoing_verification_runs(db)

    return {"status": "OK"}

@app.get("/pendingdocs/{userId}")
def get_users_pending_docs_count(userId: str, db: Session = Depends(get_db)):
    return db.query(Document).filter(Document.uploaded_by == userId, Document.status == "pending").count()

@app.get("/usertokens/{userId}")
def get_user_tokens(userId: str, db: Session = Depends(get_db)):
    user = db.query(User).filter(User.id == userId).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    return {"tokens": user.tokens}

@app.get("/files", response_model=List[Dict[str, str]])
async def get_files(db: Session = Depends(get_db)):
    """
    Az adatbázisból listázza a feltöltött dokumentumokat.
    """
    documents = db.query(Document).all()
    
    return [
        {
            "id": doc.id,
            "title": doc.title,
            "description": doc.description,
            "file_path": doc.file_path,
            "file_name": doc.file_path.split('/')[-1],
            "status": doc.status,
            "uploaded_by": doc.uploaded_by,
            "uploaded_at": doc.uploaded_at.isoformat(),
            "delete_url": f"/delete/{doc.file_path.split('/')[-1]}",
            "download_url": f"/download/{doc.file_path.split('/')[-1]}",
        }
        for doc in documents
    ]



def send_email(recipient_email: str, verification_code: str, name: str):
    sender_email = "poseidongg.noreply@gmail.com"
    sender_password = "opst qfmv gwzb lhxa"

    smtp_server = "smtp.gmail.com"
    smtp_port = 587

    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = recipient_email
    message["Subject"] = "Verifikációs kód"

    # Kép csatolása
    with open("logo.jpg", "rb") as img:
        img_data = img.read()
        image = MIMEImage(img_data)
        image.add_header("Content-ID", "<logo>")
        message.attach(image)


    # HTML template módosítva
    html_body = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta name="viewport" content="width=device-width, initial-scale=1.0" />
        <meta http-equiv="Content-Type" content="text/html; charset=UTF-8" />
        <title>Verify your email address</title>
        <style type="text/css" rel="stylesheet" media="all">
            * {{
                font-family: Arial, 'Helvetica Neue', Helvetica, sans-serif;
                box-sizing: border-box;
            }}
            body {{
                width: 100% !important;
                height: 100%;
                margin: 0;
                background-color: #000000;
                color: #ffbb00;
            }}
            a {{
                color: #ffbb00;
            }}
            .email-wrapper {{
                width: 100%;
                margin: 0;
                padding: 0;
                background-color: #000000;
            }}
            .email-content {{
                width: 100%;
                margin: 0;
                padding: 0;
            }}
            .email-masthead {{
                padding: 25px 0;
                text-align: center;
            }}
            .email-masthead_logo {{
                max-width: 200px;
                border: 0;
            }}
            .email-body {{
                width: 100%;
                margin: 0;
                padding: 0;
                border-top: 1px solid #ffbb00;
                border-bottom: 1px solid #ffbb00;
                background-color: #000000;
            }}
            .email-body_inner {{
                width: 570px;
                margin: 0 auto;
                padding: 0;
            }}
            .email-footer {{
                width: 570px;
                margin: 0 auto;
                padding: 0;
                text-align: center;
            }}
            .email-footer p {{
                color: #ffbb00;
            }}
            .content-cell {{
                padding: 35px;
            }}
            h1 {{
                margin-top: 0;
                color: #ffbb00;
                font-size: 19px;
                font-weight: bold;
                text-align: left;
            }}
            p {{
                margin-top: 0;
                color: #ffbb00;
                font-size: 16px;
                line-height: 1.5em;
                text-align: left;
            }}
            .code-box {{
                background-color: #1c1c1c;
                color: #ffbb00;
                padding: 10px;
                border-radius: 5px;
                text-align: center;
                font-size: 24px;
                font-weight: bold;
                letter-spacing: 2px;
            }}
        </style>
    </head>
    <body>
        <table class="email-wrapper" width="100%" cellpadding="0" cellspacing="0">
            <tr>
                <td align="center">
                    <table class="email-content" width="100%" cellpadding="0" cellspacing="0">
                        <!-- Logo -->
                        <tr>
                            <td class="email-masthead">
                                <img src="cid:logo" alt="Poseidon Logo" class="email-masthead_logo" />
                            </td>
                        </tr>
                        <!-- Email Body -->
                        <tr>
                            <td class="email-body" width="100%">
                                <table class="email-body_inner" align="center" width="570" cellpadding="0" cellspacing="0">
                                    <tr>
                                        <td class="content-cell">
                                            <h1>Verify your email address</h1>
                                            <p>Dear {name}!</p>
                                            <p>Thanks for signing up! Use the code below to verify your email address:</p>
                                            <div class="code-box">{verification_code}</div>
                                            <p>If you didn't sign up, you can safely ignore this email.</p>
                                            <p>Thanks,<br>The Poseidon Team</p>
                                        </td>
                                    </tr>
                                </table>
                            </td>
                        </tr>
                        <tr>
                            <td>
                                <table class="email-footer" align="center" width="570" cellpadding="0" cellspacing="0">
                                    <tr>
                                        <td class="content-cell">
                                            <p class="sub center">
                                                Poseidon<br>
                                                All rights reserved.
                                            </p>
                                        </td>
                                    </tr>
                                </table>
                            </td>
                        </tr>
                    </table>
                </td>
            </tr>
        </table>
    </body>
    </html>
    """

    message.attach(MIMEText(html_body, "html"))

    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(sender_email, sender_password)
            server.sendmail(sender_email, recipient_email, message.as_string())
        print(f"Email sikeresen elküldve {recipient_email} címre.")
    except Exception as e:
        print(f"Az email küldésének hibája: {str(e)}")


@app.get("/email/decision")
def send_email_decision(recipient_email: str, title: str, decision: str, sender: str, username: str, fileId: str, rejection_reason: str = None, db: Session = Depends(get_db)):
    sender_email = "poseidongg.noreply@gmail.com"  # A Te email címed
    sender_password = "opst qfmv gwzb lhxa"  # Az email jelszavad

    fileinfo=db.query(Document).filter(Document.id == fileId).first()
    # SMTP beállítások
    smtp_server = "smtp.gmail.com"
    smtp_port = 587  # Vagy 465 a SSL-hez

    # Email üzenet készítése
    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = recipient_email

    # Kép csatolása
    with open("logo.jpg", "rb") as img:
        img_data = img.read()
        image = MIMEImage(img_data)
        image.add_header("Content-ID", "<logo>")
        message.attach(image)

    if decision == "approved":
        message["Subject"] = "Feltöltött Fájl Jóváhagyva"
        if fileinfo.is_edit:
            html_body = f"""
            <!DOCTYPE html>
            <html lang="en">
            <head>
                <meta name="viewport" content="width=device-width, initial-scale=1.0" />
                <meta http-equiv="Content-Type" content="text/html; charset=UTF-8" />
                <title>Verify your email address</title>
                <style type="text/css" rel="stylesheet" media="all">
                    * {{
                        font-family: Arial, 'Helvetica Neue', Helvetica, sans-serif;
                        box-sizing: border-box;
                    }}
                    body {{
                        width: 100% !important;
                        height: 100%;
                        margin: 0;
                        background-color: #000000;
                        color: #ffbb00;
                    }}
                    a {{
                        color: #ffbb00;
                    }}
                    .email-wrapper {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                        background-color: #000000;
                    }}
                    .email-content {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                    }}
                    .email-masthead {{
                        padding: 25px 0;
                        text-align: center;
                    }}
                    .email-masthead_logo {{
                        max-width: 200px;
                        border: 0;
                    }}
                    .email-body {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                        border-top: 1px solid #ffbb00;
                        border-bottom: 1px solid #ffbb00;
                        background-color: #000000;
                    }}
                    .email-body_inner {{
                        width: 570px;
                        margin: 0 auto;
                        padding: 0;
                    }}
                    .email-footer {{
                        width: 570px;
                        margin: 0 auto;
                        padding: 0;
                        text-align: center;
                    }}
                    .email-footer p {{
                        color: #ffbb00;
                    }}
                    .content-cell {{
                        padding: 35px;
                    }}
                    h1 {{
                        margin-top: 0;
                        color: #ffbb00;
                        font-size: 19px;
                        font-weight: bold;
                        text-align: left;
                    }}
                    p {{
                        margin-top: 0;
                        color: #ffbb00;
                        font-size: 16px;
                        line-height: 1.5em;
                        text-align: left;
                    }}
                    .code-box {{
                        background-color: #1c1c1c;
                        color: #ffbb00;
                        padding: 10px;
                        border-radius: 5px;
                        text-align: center;
                        font-size: 24px;
                        font-weight: bold;
                        letter-spacing: 2px;
                    }}
                </style>
            </head>
            <body>
                <table class="email-wrapper" width="100%" cellpadding="0" cellspacing="0">
                    <tr>
                        <td align="center">
                            <table class="email-content" width="100%" cellpadding="0" cellspacing="0">
                                <!-- Logo -->
                                <tr>
                                    <td class="email-masthead">
                                        <img src="cid:logo" alt="Poseidon Logo" class="email-masthead_logo" />
                                    </td>
                                </tr>
                                <!-- Email Body -->
                                <tr>
                                    <td class="email-body" width="100%">
                                        <table class="email-body_inner" align="center" width="570" cellpadding="0" cellspacing="0">
                                            <tr>
                                                <td class="content-cell">
                                                    <h1>Dear {username}!</h1>
                                                    <p>Your edited file "{title}" has been approved by {sender}.</p>
                                                    <p>Thanks,<br>The Poseidon Team</p>
                                                </td>
                                            </tr>
                                        </table>
                                    </td>
                                </tr>
                                <tr>
                                    <td>
                                        <table class="email-footer" align="center" width="570" cellpadding="0" cellspacing="0">
                                            <tr>
                                                <td class="content-cell">
                                                    <p class="sub center">
                                                        Poseidon<br>
                                                        All rights reserved.
                                                    </p>
                                                </td>
                                            </tr>
                                        </table>
                                    </td>
                                </tr>
                            </table>
                        </td>
                    </tr>
                </table>
            </body>
            </html>
            """
        else:    
            html_body = f"""
            <!DOCTYPE html>
            <html lang="en">
            <head>
                <meta name="viewport" content="width=device-width, initial-scale=1.0" />
                <meta http-equiv="Content-Type" content="text/html; charset=UTF-8" />
                <title>Verify your email address</title>
                <style type="text/css" rel="stylesheet" media="all">
                    * {{
                        font-family: Arial, 'Helvetica Neue', Helvetica, sans-serif;
                        box-sizing: border-box;
                    }}
                    body {{
                        width: 100% !important;
                        height: 100%;
                        margin: 0;
                        background-color: #000000;
                        color: #ffbb00;
                    }}
                    a {{
                        color: #ffbb00;
                    }}
                    .email-wrapper {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                        background-color: #000000;
                    }}
                    .email-content {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                    }}
                    .email-masthead {{
                        padding: 25px 0;
                        text-align: center;
                    }}
                    .email-masthead_logo {{
                        max-width: 200px;
                        border: 0;
                    }}
                    .email-body {{
                        width: 100%;
                        margin: 0;
                        padding: 0;
                        border-top: 1px solid #ffbb00;
                        border-bottom: 1px solid #ffbb00;
                        background-color: #000000;
                    }}
                    .email-body_inner {{
                        width: 570px;
                        margin: 0 auto;
                        padding: 0;
                    }}
                    .email-footer {{
                        width: 570px;
                        margin: 0 auto;
                        padding: 0;
                        text-align: center;
                    }}
                    .email-footer p {{
                        color: #ffbb00;
                    }}
                    .content-cell {{
                        padding: 35px;
                    }}
                    h1 {{
                        margin-top: 0;
                        color: #ffbb00;
                        font-size: 19px;
                        font-weight: bold;
                        text-align: left;
                    }}
                    p {{
                        margin-top: 0;
                        color: #ffbb00;
                        font-size: 16px;
                        line-height: 1.5em;
                        text-align: left;
                    }}
                    .code-box {{
                        background-color: #1c1c1c;
                        color: #ffbb00;
                        padding: 10px;
                        border-radius: 5px;
                        text-align: center;
                        font-size: 24px;
                        font-weight: bold;
                        letter-spacing: 2px;
                    }}
                </style>
            </head>
            <body>
                <table class="email-wrapper" width="100%" cellpadding="0" cellspacing="0">
                    <tr>
                        <td align="center">
                            <table class="email-content" width="100%" cellpadding="0" cellspacing="0">
                                <!-- Logo -->
                                <tr>
                                    <td class="email-masthead">
                                        <img src="cid:logo" alt="Poseidon Logo" class="email-masthead_logo" />
                                    </td>
                                </tr>
                                <!-- Email Body -->
                                <tr>
                                    <td class="email-body" width="100%">
                                        <table class="email-body_inner" align="center" width="570" cellpadding="0" cellspacing="0">
                                            <tr>
                                                <td class="content-cell">
                                                    <h1>Dear {username}!</h1>
                                                    <p>Your uploaded file "{title}" has been approved by {sender}.</p>
                                                    <div class="code-box">+4 Quiz Tokens</div>
                                                    <p>You can now use these tokens for quiz generations.</p>
                                                    <p>Thanks,<br>The Poseidon Team</p>
                                                </td>
                                            </tr>
                                        </table>
                                    </td>
                                </tr>
                                <tr>
                                    <td>
                                        <table class="email-footer" align="center" width="570" cellpadding="0" cellspacing="0">
                                            <tr>
                                                <td class="content-cell">
                                                    <p class="sub center">
                                                        Poseidon<br>
                                                        All rights reserved.
                                                    </p>
                                                </td>
                                            </tr>
                                        </table>
                                    </td>
                                </tr>
                            </table>
                        </td>
                    </tr>
                </table>
            </body>
            </html>
            """
    elif decision == "rejected":
        message["Subject"] = "Feltöltött Fájl Elutasítva"
        html_body = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta name="viewport" content="width=device-width, initial-scale=1.0" />
            <meta http-equiv="Content-Type" content="text/html; charset=UTF-8" />
            <title>Verify your email address</title>
            <style type="text/css" rel="stylesheet" media="all">
                * {{
                    font-family: Arial, 'Helvetica Neue', Helvetica, sans-serif;
                    box-sizing: border-box;
                }}
                body {{
                    width: 100% !important;
                    height: 100%;
                    margin: 0;
                    background-color: #000000;
                    color: #ffbb00;
                }}
                a {{
                    color: #ffbb00;
                }}
                .email-wrapper {{
                    width: 100%;
                    margin: 0;
                    padding: 0;
                    background-color: #000000;
                }}
                .email-content {{
                    width: 100%;
                    margin: 0;
                    padding: 0;
                }}
                .email-masthead {{
                    padding: 25px 0;
                    text-align: center;
                }}
                .email-masthead_logo {{
                    max-width: 200px;
                    border: 0;
                }}
                .email-body {{
                    width: 100%;
                    margin: 0;
                    padding: 0;
                    border-top: 1px solid #ffbb00;
                    border-bottom: 1px solid #ffbb00;
                    background-color: #000000;
                }}
                .email-body_inner {{
                    width: 570px;
                    margin: 0 auto;
                    padding: 0;
                }}
                .email-footer {{
                    width: 570px;
                    margin: 0 auto;
                    padding: 0;
                    text-align: center;
                }}
                .email-footer p {{
                    color: #ffbb00;
                }}
                .content-cell {{
                    padding: 35px;
                }}
                h1 {{
                    margin-top: 0;
                    color: #ffbb00;
                    font-size: 19px;
                    font-weight: bold;
                    text-align: left;
                }}
                p {{
                    margin-top: 0;
                    color: #ffbb00;
                    font-size: 16px;
                    line-height: 1.5em;
                    text-align: left;
                }}
                .code-box {{
                    background-color: #1c1c1c;
                    color: #ffbb00;
                    padding: 10px;
                    border-radius: 5px;
                    text-align: center;
                    font-size: 24px;
                    font-weight: bold;
                    letter-spacing: 2px;
                }}
            </style>
        </head>
        <body>
            <table class="email-wrapper" width="100%" cellpadding="0" cellspacing="0">
                <tr>
                    <td align="center">
                        <table class="email-content" width="100%" cellpadding="0" cellspacing="0">
                            <!-- Logo -->
                            <tr>
                                <td class="email-masthead">
                                    <img src="cid:logo" alt="Poseidon Logo" class="email-masthead_logo" />
                                </td>
                            </tr>
                            <!-- Email Body -->
                            <tr>
                                <td class="email-body" width="100%">
                                    <table class="email-body_inner" align="center" width="570" cellpadding="0" cellspacing="0">
                                        <tr>
                                            <td class="content-cell">
                                                <h1>Dear {username}!</h1>
                                                <p>Your uploaded file "{title}" has been rejected by {sender}.</p>
                                                <p>Reason: {rejection_reason} </p>
                                                <p>Try uploading another file.</p>
                                                <p>Thanks,<br>The Poseidon Team</p>
                                            </td>
                                        </tr>
                                    </table>
                                </td>
                            </tr>
                            <tr>
                                <td>
                                    <table class="email-footer" align="center" width="570" cellpadding="0" cellspacing="0">
                                        <tr>
                                            <td class="content-cell">
                                                <p class="sub center">
                                                    Poseidon<br>
                                                    All rights reserved.
                                                </p>
                                            </td>
                                        </tr>
                                    </table>
                                </td>
                            </tr>
                        </table>
                    </td>
                </tr>
            </table>
        </body>
        </html>
        """
    else:
        return {"status": "ERROR", "message": "Invalid decision type"}

    message.attach(MIMEText(html_body, "html"))

    # SMTP kapcsolat létrehozása és az email elküldése
    try:
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()  # TLS titkosítás engedélyezése
            server.login(sender_email, sender_password)
            server.sendmail(sender_email, recipient_email, message.as_string())
        print(f"Email sikeresen elküldve {recipient_email} címre.")
    except Exception as e:
        print(f"Az email küldésének hibája: {str(e)}")

    return {"status": "OK"}



class ConfirmationResponse(BaseModel):
    status: str
    error_id: Optional[str] = None



@app.post("/confirm", response_model=ConfirmationResponse)
def confirm_verification(
    entity_type: str, 
    entity_id: str, 
    verification_process: str, 
    code: str,
    request: Request,
    service_provider_id: str = 'VB',
    session: Session = Depends(get_db)
):
    base = BaseClass()



    VERIFICATION_EXPIRE_DAYS=5000
    print(f"DEBUG: entity_id={entity_id}")

    verification_run_id = base.get_verification_run_id(entity_id, session)


    
    run = base.get_verification_run(verification_run_id, service_provider_id, entity_type, entity_id, verification_process, session)

    if run.status=='FAILED':
        return {"status": "ERROR", "error_id": "FAILED_VERIFICATION. CREATE NEW ACCOUNT"}


    proof = run.proofs[0]



    if run.status=='ONGOING':
        user_ip = request.client.host

        proof.ip_address = user_ip

        if not run or proof.verification_code != code:
            
            run.remaining_tries -= 1
            base.update_verification_status(verification_run=run, new_status="FAILED" if run.remaining_tries <= -50000000000000 else run.status, session=session)

            if run.remaining_tries <= -5000000000000:
                #base.update_proof_status(proof, "REJECTED", datetime.now(), datetime.now() + timedelta(days=365), session)
                proof.status = "REJECTED"
                run.fail_reason = "TOO_MANY_TRIES"
                session.commit()  
                session.close()
                return {"status": "ERROR", "error_id": "TOO_MANY_TRIES"}


            session.commit()  
            session.close()

            return {"status": "ERROR", "error_id": "BAD_CODE"}
        
        proof.correct_code_submission_time=datetime.now()
        base.update_verification_status(run, new_status="FINISHED", session=session)
        #base.update_proof_status(proof, "APPROVED", datetime.now(), datetime.now() + timedelta(days=365), session)
        proof.status = "APPROVED"
        proof.expirationDate = datetime.now() + timedelta(days=VERIFICATION_EXPIRE_DAYS)

        verification_data = {

            "serviceProviderID": run.serviceProviderID,
            "verificationTypeCode": run.verificationTypeCode,
            "verificationProcessCode": run.verificationProcessCode,
            "verificationRunID": run.id,
            "entityType": entity_type,
            "entityID": entity_id,
            "status": "VALID",
            "effective_date": datetime.now(),
            "expiration_date": datetime.now() + timedelta(days=VERIFICATION_EXPIRE_DAYS),
            "data": {"email": proof.main_param}
        }
        base.create_verification(verification_data, session)

        '''
        # Call PHP setLevelFromVerification API
        php_result = call_php_api_set_level(entity_id, "PHONE", True)

        if not php_result["success"]:
            return {"status": "ERROR", "error_id": "PHP_API_CALL_FAILED"}
        '''
        success=verify_get_user_from_db(email=proof.main_param, db=session)
                         
    session.commit()  
    session.close()


    return {"status": "OK"}

class CancelVerificationResponse(BaseModel):
    status: str
    error_id: str = None

@app.post("/cancel", response_model=CancelVerificationResponse)
def cancel_verification(
    entityType: str, 
    entityID: str, 
    verification_process: str, 
    verificationRunID: str,
    serviceProviderID: str = 'VB',
    session: Session = Depends(get_db)
):
    base = BaseClass()


    run = base.get_verification_run(verificationRunID, serviceProviderID, entityType, entityID, verification_process, session)
    
    if not run or run.status != "ONGOING":
        return CancelVerificationResponse(status="ERROR", error_id="INVALID_VERIFICATION")

    base.update_verification_status(verification_run=run, new_status='CANCELLED', session=session)
    base.delete_proof(run.id, session)

    session.commit()  
    session.close()

    return CancelVerificationResponse(status="OK")


def calculate_linear_wait_time(attempt_number: int, waitMinutes: int) -> timedelta:
    return timedelta(minutes=waitMinutes)

def calculate_exponential_wait_time(attempt_number: int, waitMinutes: int) -> timedelta:
    return timedelta(minutes=waitMinutes ** (attempt_number + 1))

@app.post("/resend")
def resend_code(
    entity_type: str, 
    entity_id: str, 
    verification_process: str = "EMAIL", 
    service_provider_id: str = 'VB',
    session: Session = Depends(get_db),
    method="exponential"):

    base = BaseClass()
    current_timestamp = datetime.now(timezone.utc)  


    CODE_LENGTH=6
    MAX_RETRY_PROCESS=50000000
    MAX_RETRY_PROCESS_WAIT_TIME_MINUTES=1.5
    MAX_RETRY_PROCESS_METHOD="exponential"

    try:
        session.begin()
        
        verification_run_id = base.get_verification_run_id(entity_id, session)
        
        
        print(f"DEBUG: verification_run_id={verification_run_id}")

        if verification_run_id is None:
            return {"error": "No verification run found for given entity_id"}
        
        verification_run = base.get_verification_run(
            verification_run_id, service_provider_id, entity_type, entity_id, 
            verification_process, session
        )

        if verification_run.status != "ONGOING":
            return {"error": "VERIFICATION_RUN_NOT_ONGOING"}
        
        proof = verification_run.proofs[0]
        
        if verification_run.try_count >= MAX_RETRY_PROCESS:
            return {"error": "MAX_RESEND_ATTEMPTS_REACHED"}
        
        if MAX_RETRY_PROCESS_METHOD == "linear":
            wait_time = calculate_linear_wait_time(verification_run.try_count, MAX_RETRY_PROCESS_WAIT_TIME_MINUTES)
        elif MAX_RETRY_PROCESS_METHOD == "exponential":
            wait_time = calculate_exponential_wait_time(verification_run.try_count, MAX_RETRY_PROCESS_WAIT_TIME_MINUTES)
        else:
            return {"error": "Invalid method."}
        
        #next_resend_time = verification_run.last_try_timestamp + wait_time  
        
        if verification_run.try_count != 0:
            if MAX_RETRY_PROCESS_METHOD == "linear":
                wait_time_last = calculate_linear_wait_time((verification_run.try_count)-1, MAX_RETRY_PROCESS_WAIT_TIME_MINUTES)
            elif MAX_RETRY_PROCESS_METHOD == "exponential":
                wait_time_last = calculate_exponential_wait_time((verification_run.try_count)-1, MAX_RETRY_PROCESS_WAIT_TIME_MINUTES)

            last_next_resend_time = verification_run.last_try_timestamp + wait_time_last

            if current_timestamp < last_next_resend_time:
                return {"error": f"NEXT_RESEND_AVAILABLE_AT {last_next_resend_time}"}
        
        verification_run.try_count += 1
        
        old_prefix = proof.prefix
        old_verification_code = proof.verification_code

        while True:
            new_prefix = ''.join([str(random.randint(0, 9)) for _ in range(3)])
            new_verification_code = ''.join([str(random.randint(0, 9)) for _ in range(CODE_LENGTH)])
            
            if new_prefix != old_prefix and new_verification_code != old_verification_code:
                break

        verification_run.last_try_timestamp = current_timestamp
        proof.prefix = new_prefix
        proof.verification_code = new_verification_code
        #send code to phone number


                #send email to address
        #send_email(db_user.email, verification_code)
        db_user=session.query(User).filter(User.id == entity_id).first()

        send_email(proof.main_param, new_verification_code, db_user.name)

        session.commit()





        #next_resend_time = current_timestamp + wait_time
        if verification_run.try_count >= MAX_RETRY_PROCESS:
            return {
                "prefix": new_prefix,
                "status": "MAX_RESEND_ATTEMPTS_REACHED",
            }
        return {
            "prefix": new_prefix,
            "next_resend_time": verification_run.last_try_timestamp + wait_time
        }

    finally:
        session.close()



@app.get("/categories")
def get_categories(db: Session = Depends(get_db)):
    categories = db.query(Category).all()
    category_map = {cat.id: cat for cat in categories}
    
    def build_tree(category):
        return {
            "id": category.id,
            "name": category.name,
            "children": [build_tree(sub) for sub in categories if sub.parent_id == category.id]
        }
    
    return [build_tree(cat) for cat in categories if cat.parent_id is None]


@app.get("/files/{category_id}")
def get_documents_by_category(category_id: str, db: Session = Depends(get_db)):

    if category_id == "trending":
        # Lekérdezzük a 5 legnépszerűbb dokumentumot
        documents = db.query(Document) \
            .filter(Document.status == "approved") \
            .order_by(Document.popularity.desc()) \
            .limit(5) \
            .all()
        
    elif category_id == "recent":
        # Top 5 legújabb dokumentum az 'uploaded_at' mező szerint csökkenő sorrendben
        documents = db.query(Document) \
            .filter(Document.status == "approved") \
            .order_by(Document.uploaded_at.desc()) \
            .limit(5) \
            .all()
    else:
        documents = db.query(Document).filter(Document.category_id == category_id, Document.status == "approved").all()


    return [
        {
            "id": doc.id,
            "title": doc.title,
            "description": doc.description,
            "file_path": doc.file_path,
            "status": doc.status,
            "file_name": doc.file_path.split('/')[-1],
            "uploaded_by": doc.uploaded_by,
            "category_id": doc.category_id,
            "uploaded_at": doc.uploaded_at.isoformat(),
            "delete_url": f"/delete/{doc.file_path.split('/')[-1]}",
            "download_url": f"/download/{doc.file_path.split('/')[-1]}",
            "uploaded_at_display": datetime.strptime(doc.uploaded_at.isoformat(), "%Y-%m-%dT%H:%M:%S.%f").strftime("%-m/%-d/%Y"),
        }
        for doc in documents
    ]

@app.get("/filesinfo/{fileId}")
def get_documents_information(fileId: str, db: Session = Depends(get_db)):
    doc = db.query(Document).filter(Document.id == fileId).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    
    usr = db.query(User).filter(User.id == doc.uploaded_by).first()
    if not usr:
        raise HTTPException(status_code=404, detail="User not found")
    
    return {
            "response": "OK",
            "title": doc.title,
            "delete_url": f"/delete/{doc.file_path.split('/')[-1]}",
            "usremail": usr.email,
            "usrname": usr.name,
            "status": doc.status,
        }


@app.get("/moderations/files")
def get_pending_files(db: Session = Depends(get_db)):
    documents = db.query(Document).filter(Document.status == "pending").all()
    return [
        {
            "id": doc.id,
            "title": doc.title,
            "description": doc.description,
            "file_path": doc.file_path,
            "status": doc.status,
            "file_name": doc.file_path.split('/')[-1],
            "category_id": doc.category_id,
            "uploaded_by": doc.uploaded_by,
            "uploaded_at": doc.uploaded_at.isoformat(),
            "delete_url": f"/delete/{doc.file_path.split('/')[-1]}",
            "download_url": f"/download/{doc.file_path.split('/')[-1]}",
        }
        for doc in documents
    ]


@app.put("/moderations/approve/{file_id}")
async def approve_file(file_id: str, db: Session = Depends(get_db), current_user = Depends(get_current_user)):
    doc = db.query(Document).filter(Document.id == file_id, Document.status == 'pending').first()
    if not doc:
        raise HTTPException(status_code=404, detail="File not found or already processed")
    
    doc.status = 'approved'

    if not doc.is_edit:
        user = db.query(User).filter(User.id == doc.uploaded_by).first()
        if user:
            user.tokens += 4  # +4 token ha nem szerkesztett fájl
            db.commit()
            
    log = ModerationLog(
        document_id=file_id,
        moderator_id=current_user.id,
        decision='approved',
        reason=None  # Jóváhagyás esetén nincs ok
    )
    db.add(log)
    db.commit()
    return {"message": "File approved successfully"}



@app.put("/moderations/reject/{file_id}")
async def reject_file(file_id: str, reason: str, db: Session = Depends(get_db), current_user = Depends(get_current_user)):
    doc = db.query(Document).filter(Document.id == file_id, Document.status == 'pending').first()
    if not doc:
        raise HTTPException(status_code=404, detail="File not found or already processed")
    
    doc.status = 'rejected'

    log = ModerationLog(
        document_id=file_id,
        moderator_id=current_user.id,
        decision='rejected',
        reason=reason  # Elutasítás oka
    )
    db.add(log)
    #file.rejection_reason = reason
    db.commit()
    return {"message": "File rejected successfully"}


@app.post("/api/documents/{document_id}/increase_popularity")
def increase_popularity(document_id: str, db: Session = Depends(get_db)):
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")

    document.popularity += 1
    db.commit()
    db.refresh(document)

    return {"message": "Popularity increased", "new_popularity": document.popularity}



def split_text_into_chunks(text, max_length=3000):
    """ Mondatokat figyelembe véve darabolja a szöveget max_length karakteres blokkokra. """
    sentences = re.split(r'(?<=[.!?])\s+', text)  # Mondathatárok figyelembevétele
    chunks = []
    current_chunk = ""

    for sentence in sentences:
        if len(current_chunk) + len(sentence) < max_length:
            current_chunk += " " + sentence if current_chunk else sentence
        else:
            chunks.append(current_chunk)
            current_chunk = sentence

    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

def translate_large_text(text, source_lang, target_lang, max_length=3000):
    """ Nagyobb szövegek darabolva történő fordítása. """
    print("nagyszöveg")
    chunks = split_text_into_chunks(text, max_length)
    translated_chunks = [GoogleTranslator(source=source_lang, target=target_lang).translate(chunk) for chunk in chunks]
    return " ".join(translated_chunks)


def translate_text(text, source_lang, target_lang):
    return GoogleTranslator(source=source_lang, target=target_lang).translate(text)

def find_closest_word(word, word_list):
    closest_matches = difflib.get_close_matches(word, word_list, n=1)
    return closest_matches[0] if closest_matches else word

load_dotenv() 
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

class QuizRequest(BaseModel):
    text: str
    num_questions: int
    language: str  # "hu" vagy "en"

def count_tokens(text: str) -> int:
    """Durva becslés tokenek számolására (jobb, ha tiktoken könyvtárat használod)."""
    return len(text.split()) * 1.3  # Átlagosan 1 szó ~1.3 token


def generate_quiz(text, lang, max_questions):
    MAX_TOKENS = 16384

    if lang not in ["magyar", "angol"]:
        raise HTTPException(status_code=400, detail="Unsupported language. Use 'magyar' or 'angol'.")
    
    if count_tokens(text) > MAX_TOKENS:
        raise HTTPException(
            status_code=400,
            detail="A megadott szöveg túl hosszú! Próbálj meg rövidebb szöveget használni."
        )
    

    
    # Prompt készítése a választott nyelv alapján
    if lang == "magyar":
        prompt = f"""
        Kérlek, készíts {max_questions} darab feleletválasztós kvízkérdést az alábbi szövegből:
        {text}
        
        Az eredményt a következő JSON formátumban add vissza:
        {{"questions": [
            {{"question_statement": "<kérdés>", "options": ["<válasz1>", "<válasz2>", "<válasz3>", "<válasz4>"], "answer": "<helyes_válasz>"}},
            ...
        ]}}
        """
    else:
        prompt = f"""
        Generate {max_questions} multiple-choice quiz questions from the following text:
        {text}
        
        Provide the output in the following JSON format:
        {{"questions": [
            {{"question_statement": "<question>", "options": ["<option1>", "<option2>", "<option3>", "<option4>"], "answer": "<correct_option>"}},
            ...
        ]}}
        """
    
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": prompt}]
    )

    quiz_data = json.loads(response.choices[0].message.content)

    
    return quiz_data

from pdfminer.high_level import extract_text
import pdfplumber
from pptx import Presentation
import io

from docx import Document as DocxDocument

@app.get("/generate-quiz/{document_id_form}-{filename}")
async def generate_quiz_from_s3(
    filename: str,
    document_id_form: str,
    background_tasks: BackgroundTasks,
    lang: str = 'angol',
    max_questions: int = 5,
    user_id: str = 'default_user',
    db: Session = Depends(get_db)
):
    user = db.query(User).filter(User.id == user_id).first()
    if user.tokens <= 0:
        return JSONResponse(
            content={"message": "Elfogytak a kvízgeneráláshoz szükséges tokenjeid. Tölts fel fájlokat, hogy tokent szerezz!"},
            status_code=400
        )
    
    # Token levonása
    user.tokens -= 1
    db.commit()


    s3 = boto3.client('s3', aws_access_key_id=AWS_ACCESS_KEY_ID, aws_secret_access_key=AWS_SECRET_ACCESS_KEY, region_name=AWS_REGION_NAME)
    file_obj = s3.get_object(Bucket=S3_BUCKET_NAME, Key=filename)
    file_content = file_obj['Body'].read()
    file_extension = filename.split('.')[-1].lower()
    extracted_text = ""

    if file_extension == 'txt':
        # A szöveg kinyerése txt fájlból
        extracted_text = file_content.decode('utf-8')

    elif file_extension == 'docx':
        import io
        # A szöveg kinyerése docx fájlból
        doc = DocxDocument(io.BytesIO(file_content))
        extracted_text = '\n'.join([para.text for para in doc.paragraphs])

    elif file_extension == 'pdf':
        import io

        # A szöveg kinyerése pdf fájlból
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            extracted_text = '\n'.join([page.extract_text() for page in pdf.pages])

    elif file_extension == 'ppt' or file_extension == 'pptx':
        import io

        # A szöveg kinyerése ppt fájlból
        prs = Presentation(io.BytesIO(file_content))
        extracted_text = '\n'.join([slide.shapes[0].text for slide in prs.slides if hasattr(slide.shapes[0], "text")])

    else:
        return JSONResponse(content={"message": "Nem támogatott fájlformátum"}, status_code=400)

    # Kvíz generálása


    quiz_id = f"quiz_{uuid.uuid4()}"
    new_quiz = Quiz(
            id=quiz_id,
            document_id=document_id_form,  
            created_by=user_id,
            is_ready=False,
            created_at=datetime.utcnow()
        )
    db.add(new_quiz) 
    db.commit()
    db.refresh(new_quiz)

    if lang not in ["magyar", "angol"]:
        return JSONResponse(content={"message": "Nem támogatott nyelv"}, status_code=400)
    
    MAX_TOKENS = 16384


    
    if count_tokens(extracted_text) > MAX_TOKENS:
        raise HTTPException(
            status_code=400,
            detail="A megadott szöveg túl hosszú! Próbálj meg rövidebb szöveget használni."
        )
    
    background_tasks.add_task(generate_quiz_background, extracted_text, lang, max_questions, document_id_form, user_id, new_quiz.id, db)
    
    return JSONResponse(content={"message": "Kvíz generálása folyamatban...", "quiz_id": new_quiz.id})    



@app.post("/start_verification")
def start_verification(
    entity_id: str, 
    db: Session = Depends(get_db)
):
        base = BaseClass()
        db_user=db.query(User).filter(User.id == entity_id).first()
        run_duplicate=base.is_run_duplicate(entity_id=db_user.id, verification_process="EMAIL", session=db)

        if run_duplicate!="":
            new_duplicate_run = VerificationRunDuplicate(
                id=f"verification_verificationrunduplicate_{uuid.uuid4()}",
                serviceProviderID="VB",
                verificationTypeCode="EMAIL",
                entityType=db_user.role,
                entityID=db_user.id,
                verificationProcessCode="EMAIL",
                originalVerificationRunID=run_duplicate
            )
            base.create_verification_run_duplicate(new_duplicate_run)
            return {"status": "DUPLICATE_RUN_FOUND"}
        
        VERIFICATION_EXPIRE_DAYS=5000
        CODE_LENGTH=6
        TRY_EXPIRE_HOURS=24
        MAX_RRETRY_PROCESS=3
        MAX_RETRY_PROCESS_WAIT_TIME_MINUTES=3
        MAX_RETRY_PROCESS_METHOD="EXPONENTIAL"
        MAX_RETRY=3

        new_run = VerificationRun(
                id=f"verification_verificationrun_{uuid.uuid4()}",
                serviceProviderID="VB",
                verificationProcessCode="EMAIL",
                entityType=db_user.role,
                entityID=db_user.id,
                verificationTypeCode="EMAIL",
                status="ONGOING",
                vendor_status="PENDING",
                fail_reason="",
                try_count=0,
                effective_date=datetime.now(),
                expiration_date=datetime.now() + timedelta(hours=TRY_EXPIRE_HOURS),
                remaining_tries=MAX_RETRY
            )

        created_run = base.create_verification_run(new_run, session=db)

        prefix = ''.join([str(random.randint(0, 9)) for _ in range(3)])

        verification_code = ''.join([str(random.randint(0, 9)) for _ in range(CODE_LENGTH)])

        new_proof = EmailProof(
            id=f"verification_proof_{uuid.uuid4()}",
            verificationRunID=created_run.id,
            main_param=db_user.email,
            verification_code=verification_code,
            uploadDate=datetime.now(),
            expirationDate=datetime.now() + timedelta(hours=TRY_EXPIRE_HOURS),
            entityType=db_user.role,
            entityID=db_user.id,
            prefix=prefix,
            ip_address="",
            correct_code_submission_time=None,
            status="PENDING"
        )
        created_proof=base.create_proof(new_proof, session=db)

        proof = created_run.proofs[0]

        phoneresult=base.email_duplicate_check(db_user.id, email=db_user.id, session=db)
        if phoneresult=="":
            #base.update_proof_status(proof=proof, new_status="PENDING", main_param=normalized_phone, session=session)
            proof.status="PENDING"
            proof.main_param=db_user.email





        #send email to address
        #send_email(db_user.email, verification_code)
        verification_code = prefix + "-" +verification_code
        send_email(db_user.email, verification_code, db_user.name)
        "EMAIL SENT"



@app.get("/is_verified")
def is_verified(
    entity_id: str, 
    db: Session = Depends(get_db)

):
    base = BaseClass()

    is_valid = base.is_verified(
        entity_id, 
        db
    )

    is_ongoing = base.get_verification_run_two(entity_id, db)


    return {"is_verified": is_valid, "is_ongoing": is_ongoing}


@app.get("/get-quiz/{quiz_id}")
async def get_quiz(quiz_id: str, db: Session = Depends(get_db)):
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Kvíz nem található")

    questions = db.query(Question).filter(Question.quiz_id == quiz_id).all()
    quiz_data = {"questions": []}
    for question in questions:
        answers = db.query(Answer).filter(Answer.question_id == question.id).all()
        options = [answer.answer_text for answer in answers]
        correct = next((answer.answer_text for answer in answers if answer.is_correct), None)
        quiz_data["questions"].append({
            "question": question.question_text,
            "options": options,
            "correct": correct
        })
    return quiz_data


@app.post("/save-quiz-result")
async def save_quiz_result(quiz_id: str, score: int, user_id: str, db: Session = Depends(get_db)):
    # Ellenőrizd, hogy létezik-e a kvíz
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if not quiz:
        raise HTTPException(status_code=404, detail="Kvíz nem található")

    # Eredmény mentése
    new_result = QuizResult(
        quiz_id=quiz_id,
        user_id=user_id,
        score=score
    )
    db.add(new_result)
    db.commit()
    db.refresh(new_result)
    return {"message": "Eredmény sikeresen elmentve!"}



@app.get("/quiz-category")
async def get_quiz_category(quiz_id: str, db: Session = Depends(get_db)):
    # Csatlakozó lekérdezés a kategória nevéhez
    category_name = (
        db.query(Category.name)
        .join(Document, Document.category_id == Category.id)
        .join(Quiz, Quiz.document_id == Document.id)
        .filter(Quiz.id == quiz_id)
        .first()
    )

    if not category_name:
        raise HTTPException(status_code=404, detail="Kategória nem található")

    return {"category_name": category_name[0]}


@app.get("/quiz-results")
async def get_quiz_results(
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    # Lekérjük a bejelentkezett felhasználóhoz tartozó eredményeket
    results = db.query(QuizResult).filter(QuizResult.user_id == current_user.id).all()
    output = []
    for res in results:
        # Lekérjük a Quiz rekordot
        quiz = db.query(Quiz).filter(Quiz.id == res.quiz_id).first()
        if not quiz:
            continue
        # Lekérjük a hozzá tartozó Document-et
        document = db.query(Document).filter(Document.id == quiz.document_id).first()
        if not document:
            continue
        # Kérdések számának lekérése a Quiz-hez
        total_questions = db.query(Question).filter(Question.quiz_id == quiz.id).count()
        # A kategória lekérése a Document alapján
        category = db.query(Category).filter(Category.id == document.category_id).first()
        category_name = category.name if category else "Ismeretlen kategória"
        output.append({
            "quiz_result_id": res.id,
            "quiz_id": res.quiz_id,
            "score": res.score,
            "total_questions": total_questions,
            "category": category_name,
            "document_name": document.title,
            "completed_at": res.completed_at
        })
    return output

@app.get("/quiz-all")
async def get_all_quizzes(db: Session = Depends(get_db)):
    # Az összes kvíz lekérése
    quizzes = db.query(Quiz).all()
    output = []
    
    for quiz in quizzes:
        # Lekérjük a hozzá tartozó Document-et
        document = db.query(Document).filter(Document.id == quiz.document_id).first()
        if not document:
            continue
        
        # Kérdések számának lekérése a Quiz-hez
        total_questions = db.query(Question).filter(Question.quiz_id == quiz.id).count()
        
        # A kategória lekérése a Document alapján
        category = db.query(Category).filter(Category.id == document.category_id).first()
        category_name = category.name if category else "Ismeretlen kategória"
        
        output.append({
            "quiz_id": quiz.id,
            "total_questions": total_questions,
            "category": category_name,
            "document_name": document.title,
            "created_at": quiz.created_at
        })
    
    return output


@app.delete("/delete-quiz-result", status_code=status.HTTP_200_OK)
async def delete_quiz_result(
    quiz_result_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    # Csak a bejelentkezett felhasználó eredményét töröljük
    quiz_result = db.query(QuizResult).filter(
        QuizResult.id == quiz_result_id,
        QuizResult.user_id == current_user.id
    ).first()
    if not quiz_result:
        raise HTTPException(status_code=404, detail="Kvíz eredmény nem található")
    db.delete(quiz_result)
    db.commit()
    return {"message": "Kvíz eredmény törölve!"}



def generate_quiz_background(extracted_text, lang, max_questions, document_id_form, user_id, quiz_id, db):
        quiz_data = generate_quiz(extracted_text, lang, max_questions)

        # **1. Létrehozzuk a kvízt az adatbázisban**
        existing_quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
        if not existing_quiz:
            print(f"Hiba: Nincs ilyen kvíz ID: {quiz_id}")
            return  # Ha nincs ilyen kvíz, kilépünk

        # 🔄 Frissítjük a meglévő `Quiz` rekordot
        existing_quiz.created_by = user_id
        existing_quiz.created_at = datetime.utcnow()
        db.commit()

            # **2. Hozzáadjuk a kérdéseket**
        for question in quiz_data["questions"]:
            question_id = f"question_{uuid.uuid4()}"
            new_question = Question(
                id=question_id,
                quiz_id=quiz_id,
                question_text=question["question_statement"],
                correct_answer=question["answer"]
            )
            db.add(new_question)
            db.commit()  # Commitáljuk a kérdést, hogy biztosítsuk, hogy az id már létezik

            # **3. Hozzáadjuk a válaszokat**
            options = question["options"]  # Helyes válasz is bekerül
            for option in options:
                new_answer = Answer(
                    id=f"answer_{uuid.uuid4()}",
                    question_id=question_id,  # Mivel commitáltuk a kérdést, biztosak vagyunk benne, hogy az id létezik
                    answer_text=option,
                    is_correct=(option == question["answer"])
                )
                db.add(new_answer)
        existing_quiz.is_ready = True
        db.commit()  # Az összes változtatást véglegesítjük
        return JSONResponse(content={"message": "Kvíz elmentve!", "quiz_id": quiz_id})



@app.get("/check-quiz-status/{quiz_id}")
async def check_quiz_status(quiz_id: str, db: Session = Depends(get_db)):
    quiz = db.query(Quiz).filter(Quiz.id == quiz_id).first()
    if quiz and quiz.is_ready:
        return {"ready": True}
    return {"ready": False}


class ModerationLogResponse(BaseModel):
    id: str
    document_id: str
    document_title: str  # Hozzáadjuk a dokumentum címét
    moderator_id: str
    email: str  # Hozzáadjuk a moderátor e-mail címét
    decision: str
    reason: Optional[str] = None
    created_at: datetime

    class Config:
        orm_mode = True


@app.get("/moderation-logs", response_model=List[ModerationLogResponse])
def get_recent_moderation_logs(db: Session = Depends(get_db)):
    # Kérjük le az 5 legutóbbi moderációs logot, csatolva a dokumentumot és a moderátort
    recent_logs = (
        db.query(ModerationLog, Document.title, User.email)
        .join(Document, ModerationLog.document_id == Document.id)  # Csatlakozás a Document táblához
        .join(User, Document.uploaded_by == User.id)  # Csatlakozás a User táblához
        .order_by(ModerationLog.created_at.desc())
        .limit(5)
        .all()
    )

    # Az eredmény átalakítása a válasz modellhez
    result = []
    for log, doc_title, user_email in recent_logs:
        result.append(
            ModerationLogResponse(
                id=log.id,
                document_id=log.document_id,
                document_title=doc_title,
                moderator_id=log.moderator_id,
                email=user_email,
                decision=log.decision,
                reason=log.reason,
                created_at=log.created_at,
            )
        )

    return result